"""
Astra AI - File Manager
Reads, writes, and manages files including documents (PDF, DOCX, XLSX, PPTX).
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
from datetime import datetime, timezone
from loguru import logger


class FileManager:
    """
    File management system supporting:
    - File read/write operations
    - PDF parsing
    - Word document (DOCX) parsing
    - Excel spreadsheet (XLSX) parsing
    - PowerPoint (PPTX) parsing
    - File search
    """

    def __init__(self):
        self._initialized = False
        self.supported_extensions = {
            ".txt": "text", ".md": "markdown", ".py": "python", ".js": "javascript",
            ".ts": "typescript", ".html": "html", ".css": "css", ".json": "json",
            ".yaml": "yaml", ".yml": "yaml", ".xml": "xml", ".csv": "csv",
            ".pdf": "pdf", ".docx": "word", ".xlsx": "excel", ".pptx": "powerpoint",
            ".jpg": "image", ".jpeg": "image", ".png": "image", ".gif": "image", ".svg": "image",
        }

    async def initialize(self):
        logger.info("Initializing File Manager...")
        self._initialized = True
        logger.info("File Manager initialized")

    async def read_file(self, path: str) -> Dict[str, Any]:
        file_path = Path(path)
        if not file_path.exists():
            return {"success": False, "error": f"File not found: {path}"}
        ext = file_path.suffix.lower()
        file_type = self.supported_extensions.get(ext, "binary")

        try:
            if file_type in ("text", "markdown", "python", "javascript", "typescript", "html", "css", "json", "yaml", "xml", "csv"):
                return await self._read_text_file(file_path)
            elif file_type == "pdf":
                return await self._read_pdf(file_path)
            elif file_type == "word":
                return await self._read_docx(file_path)
            elif file_type == "excel":
                return await self._read_xlsx(file_path)
            elif file_type == "powerpoint":
                return await self._read_pptx(file_path)
            elif file_type == "image":
                return await self._read_image(file_path)
            else:
                return await self._read_binary(file_path)
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def _read_text_file(self, file_path: Path) -> Dict[str, Any]:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        return {"success": True, "content": content, "file_type": "text", "size": file_path.stat().st_size, "filename": file_path.name, "path": str(file_path)}

    async def _read_pdf(self, file_path: Path) -> Dict[str, Any]:
        try:
            import pypdfium2 as pdfium
            pdf = pdfium.PdfDocument(str(file_path))
            text_parts = []
            for i in range(len(pdf)):
                page = pdf[i]
                text = page.get_text_bounded()
                text_parts.append(f"--- Page {i+1} ---\n{text}")
            return {"success": True, "content": "\n\n".join(text_parts), "file_type": "pdf", "pages": len(pdf), "filename": file_path.name, "path": str(file_path)}
        except ImportError:
            return {"success": False, "error": "PDF support requires pypdfium2"}

    async def _read_docx(self, file_path: Path) -> Dict[str, Any]:
        try:
            from docx import Document
            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs]
            content = "\n".join(paragraphs)
            tables = []
            for table in doc.tables:
                for row in table.rows:
                    tables.append(" | ".join(cell.text for cell in row.cells))
            if tables:
                content += "\n\n--- Tables ---\n" + "\n".join(tables)
            return {"success": True, "content": content, "file_type": "docx", "paragraphs": len(paragraphs), "tables": len(tables), "filename": file_path.name, "path": str(file_path)}
        except ImportError:
            return {"success": False, "error": "DOCX support requires python-docx"}

    async def _read_xlsx(self, file_path: Path) -> Dict[str, Any]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets_data = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                sheets_data.append({"name": sheet_name, "rows": len(rows), "columns": len(rows[0]) if rows else 0,
                                    "data": "\n".join("\t".join(r) for r in rows[:100])})
            return {"success": True, "content": "\n\n".join(f"--- Sheet: {s['name']} ({s['rows']}x{s['columns']}) ---\n{s['data']}" for s in sheets_data),
                    "file_type": "xlsx", "sheets": sheets_data, "filename": file_path.name, "path": str(file_path)}
        except ImportError:
            return {"success": False, "error": "XLSX support requires openpyxl"}

    async def _read_pptx(self, file_path: Path) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            return {"success": True, "content": "\n\n".join(slides_text), "file_type": "pptx", "slides": len(prs.slides), "filename": file_path.name, "path": str(file_path)}
        except ImportError:
            return {"success": False, "error": "PPTX support requires python-pptx"}

    async def _read_image(self, file_path: Path) -> Dict[str, Any]:
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                return {"success": True, "content": f"[Image: {file_path.name}, {img.size[0]}x{img.size[1]}, {img.format}]",
                        "file_type": file_path.suffix[1:], "width": img.size[0], "height": img.size[1], "format": img.format,
                        "filename": file_path.name, "path": str(file_path)}
        except ImportError:
            return {"success": False, "error": "Image support requires Pillow"}

    async def _read_binary(self, file_path: Path) -> Dict[str, Any]:
        size = file_path.stat().st_size
        return {"success": True, "content": f"[Binary file: {file_path.name}, {size} bytes]", "file_type": "binary", "size": size, "filename": file_path.name, "path": str(file_path)}

    async def write_file(self, path: str, content: str) -> Dict[str, Any]:
        try:
            file_path = Path(path)
            file_path.parent.mkdir(parents=True, exist_ok=True)
            file_path.write_text(content, encoding="utf-8")
            return {"success": True, "path": str(file_path), "size": len(content)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def list_directory(self, path: str) -> Dict[str, Any]:
        try:
            dir_path = Path(path)
            if not dir_path.exists() or not dir_path.is_dir():
                return {"success": False, "error": f"Directory not found: {path}"}
            items = []
            for item in sorted(dir_path.iterdir()):
                items.append({"name": item.name, "type": "directory" if item.is_dir() else "file",
                              "size": item.stat().st_size if item.is_file() else 0,
                              "modified": datetime.fromtimestamp(item.stat().st_mtime, tz=timezone.utc).isoformat()})
            return {"success": True, "path": path, "items": items}
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def search_files(self, pattern: str, path: Optional[str] = None) -> Dict[str, Any]:
        try:
            search_path = Path(path) if path else Path.cwd()
            matches = list(search_path.rglob(pattern))
            return {"success": True, "pattern": pattern, "matches": [str(m.relative_to(search_path)) for m in matches[:100]], "total": len(matches)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_status(self) -> Dict[str, Any]:
        return {"initialized": self._initialized, "supported_extensions": len(self.supported_extensions)}
